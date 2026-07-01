"""
文档加载器 —— 解析 PDF / TXT / MD / DOCX / XLSX / PPTX / 图片
"""
from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


# ====== OCR Lazy Loader ======
_ocr_engine = None

def _get_ocr():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR
        _ocr_engine = RapidOCR()
    return _ocr_engine


def load_pdf(path: str | Path) -> str:
    """读取 PDF 文件，返回全文"""
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def load_txt(path: str | Path) -> str:
    """读取纯文本文件（自动检测 UTF-8 / GBK）"""
    raw = Path(path).read_bytes()
    for enc in ("utf-8", "gbk", "gb18030", "utf-16"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    # 最后的 fallback：忽略无法解码的字符
    return raw.decode("utf-8", errors="ignore")


def load_md(path: str | Path) -> str:
    """读取 Markdown 文件（和 txt 一样处理）"""
    return Path(path).read_text(encoding="utf-8")


def load_docx(path: str | Path) -> str:
    """读取 Word 文档"""
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def load_xlsx(path: str | Path) -> str:
    """读取 Excel 文件，每个单元格内容用制表符/换行分隔"""
    wb = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"【Sheet: {sheet}】")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells)
            if line.strip():
                lines.append(line)
    return "\n".join(lines)


def load_pptx(path: str | Path) -> str:
    """读取 PowerPoint 文件"""
    prs = Presentation(str(path))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            lines.append(f"【Slide {slide_num}】")
            lines.extend(texts)
    return "\n".join(lines)


def load_image(path: str | Path) -> str:
    """OCR 识别图片中的文字"""
    import numpy as np
    from PIL import Image

    img = Image.open(path).convert("RGB")
    ocr = _get_ocr()
    result, _ = ocr(np.array(img))
    if not result:
        return ""
    lines = [item[1] for item in result]
    return "\n".join(lines)


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}


def load_document(path: str | Path) -> str:
    """根据扩展名自动选择解析器"""
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return load_pdf(path)
    elif suffix in (".md", ".markdown"):
        return load_md(path)
    elif suffix == ".docx":
        return load_docx(path)
    elif suffix in (".xlsx", ".xls"):
        return load_xlsx(path)
    elif suffix == ".pptx":
        return load_pptx(path)
    elif suffix in IMAGE_SUFFIXES:
        return load_image(path)
    else:
        return load_txt(path)
